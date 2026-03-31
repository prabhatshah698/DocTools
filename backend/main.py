from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

from docx import Document
from reportlab.pdfgen import canvas
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
import qrcode

import os
import uuid
import time
import shutil
from threading import Thread

app = FastAPI()

# ---------------- FIX: REMOVE OLD CONFLICT FOLDER ----------------
if os.path.exists("pptx"):
    shutil.rmtree("pptx")

# ---------------- CORS ----------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------- FOLDERS ----------------
UPLOAD_DIR = "uploads"
PDF_DIR = "pdfs"
DOCX_DIR = "docx"
PPTX_DIR = "pptx_files"   # ✅ FIXED
QR_DIR = "qr_codes"

for folder in [UPLOAD_DIR, PDF_DIR, DOCX_DIR, PPTX_DIR, QR_DIR]:
    os.makedirs(folder, exist_ok=True)

# ---------------- TTL ----------------
FILE_TTL = 30 * 60


# ---------------- CLEANUP ----------------
def cleanup_files():
    while True:
        now = time.time()
        folders = [UPLOAD_DIR, PDF_DIR, DOCX_DIR, PPTX_DIR, QR_DIR]

        for folder in folders:
            for filename in os.listdir(folder):
                file_path = os.path.join(folder, filename)

                if os.path.isfile(file_path):
                    if now - os.path.getmtime(file_path) > FILE_TTL:
                        try:
                            os.remove(file_path)
                        except:
                            pass

        time.sleep(300)


@app.on_event("startup")
def start_cleanup():
    Thread(target=cleanup_files, daemon=True).start()


# ---------------- ROOT ----------------
@app.get("/")
def home():
    return {"message": "Backend running 🚀"}


# =====================================================
# WORD → PDF
# =====================================================
@app.post("/word-to-pdf/")
async def word_to_pdf(file: UploadFile = File(...)):
    if not file.filename.endswith(".docx"):
        raise HTTPException(400, "Only DOCX allowed")

    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}.docx")
    output_path = os.path.join(PDF_DIR, f"{file_id}.pdf")

    with open(input_path, "wb") as f:
        f.write(await file.read())

    try:
        doc = Document(input_path)
        text = "\n".join([p.text for p in doc.paragraphs])

        c = canvas.Canvas(output_path)
        y = 800

        for line in text.split("\n"):
            c.drawString(50, y, line[:100])
            y -= 20
            if y < 50:
                c.showPage()
                y = 800

        c.save()

    except Exception as e:
        raise HTTPException(500, str(e))

    return FileResponse(output_path, filename="output.pdf")


# =====================================================
# PDF → WORD
# =====================================================
@app.post("/pdf-to-word/")
async def pdf_to_word(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}.pdf")
    output_path = os.path.join(DOCX_DIR, f"{file_id}.docx")

    with open(input_path, "wb") as f:
        f.write(await file.read())

    try:
        doc = Document()

        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    doc.add_paragraph(text)

        doc.save(output_path)

    except Exception as e:
        raise HTTPException(500, str(e))

    return FileResponse(output_path, filename="output.docx")


# =====================================================
# PDF → PPT
# =====================================================
@app.post("/pdf-to-ppt/")
async def pdf_to_ppt(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}.pdf")
    output_path = os.path.join(PPTX_DIR, f"{file_id}.pptx")

    with open(input_path, "wb") as f:
        f.write(await file.read())

    try:
        prs = Presentation()

        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    box = slide.shapes.add_textbox(
                        Inches(1), Inches(1), Inches(8), Inches(5)
                    )
                    tf = box.text_frame

                    for line in text.split("\n"):
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(16)

        prs.save(output_path)

    except Exception as e:
        raise HTTPException(500, str(e))

    return FileResponse(output_path, filename="output.pptx")


# =====================================================
# WORD → PPT
# =====================================================
@app.post("/word-to-ppt/")
async def word_to_ppt(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}.docx")
    output_path = os.path.join(PPTX_DIR, f"{file_id}.pptx")

    with open(input_path, "wb") as f:
        f.write(await file.read())

    try:
        doc = Document(input_path)
        prs = Presentation()

        for para in doc.paragraphs:
            if para.text.strip():
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Slide"
                slide.placeholders[1].text = para.text

        prs.save(output_path)

    except Exception as e:
        raise HTTPException(500, str(e))

    return FileResponse(output_path, filename="output.pptx")


# =====================================================
# PPT → WORD
# =====================================================
@app.post("/ppt-to-word/")
async def ppt_to_word(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}.pptx")
    output_path = os.path.join(DOCX_DIR, f"{file_id}.docx")

    with open(input_path, "wb") as f:
        f.write(await file.read())

    try:
        prs = Presentation(input_path)
        doc = Document()

        for i, slide in enumerate(prs.slides):
            doc.add_heading(f"Slide {i+1}", level=1)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    doc.add_paragraph(shape.text)

        doc.save(output_path)

    except Exception as e:
        raise HTTPException(500, str(e))

    return FileResponse(output_path, filename="output.docx")


# =====================================================
# QR CODE
# =====================================================
@app.post("/generate-qr/")
async def generate_qr(data: str = Form(...)):
    file_id = str(uuid.uuid4())
    output_path = os.path.join(QR_DIR, f"{file_id}.png")

    img = qrcode.make(data)
    img.save(output_path)

    return FileResponse(output_path, filename="qrcode.png")